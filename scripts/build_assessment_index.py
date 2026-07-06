"""Build a lightweight assessment criteria index from the local College tree.

The public app should not ship raw college PDFs, slide decks, or personal
assessment packs. This script extracts criteria-like lines and source metadata
from local AI Programming files, then writes a compact JS module used by the
browser demo.
"""

from __future__ import annotations

import argparse
import json
import re
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Callable

from docx import Document
from pptx import Presentation

try:
    from pypdf import PdfReader
except Exception:  # pragma: no cover - optional local dependency
    PdfReader = None


SKIP_PARTS = {
    ".git",
    ".venv",
    "__pycache__",
    "dist",
    "node_modules",
    "site-packages",
    "Extracurricular\\ALI Assessment Learning Inspector",
}

SUPPORTED_EXTENSIONS = {".txt", ".md", ".docx", ".pptx", ".pdf", ".ipynb"}
TEXT_EXTENSIONS = {".txt", ".md", ".ipynb"}

RELEVANT_PATH = re.compile(
    r"(assessment|criteria|unit|lesson|evidence|support|slide|vle|coursework|assignment|project)",
    re.I,
)
CRITERION_START = re.compile(r"^(?:[-*•]\s*)?(?:AC\s*)?\d+\.\d+\s*(?:[-:|.)]\s*)?.+", re.I)
ASSESSMENT_CONTEXT = re.compile(r"(assessment criteria covered|covers assessment criteria|assessment criteria:-)", re.I)
CRITERION_VERB = re.compile(
    r"\b(explain|outline|describe|demonstrate|create|use|apply|implement|write|run|save|load|organise|check|merge|mark)\b",
    re.I,
)
PRIVATE_OR_EVIDENCE_LINE = re.compile(
    r"(C:\\|https?://|Rob\s+Johnston|screenshot|evidence table|local evidence|template source|"
    r"document type|url evidence|learner details|assessment name|github evidence folder|"
    r"video evidence folder|voiceover|file type:|there are \d+ screenshot|\|\s*\||\.(png|mp4|mp3)|"
    r"\b(this shows|visible|opened|button|shown|starts the|confirms|proves|homepage|screen)\b)",
    re.I,
)
NOISE_LINE = re.compile(r"^(page\s+\d+|\d+|copyright|all rights reserved)$", re.I)


@dataclass(frozen=True)
class SourceDoc:
    path: Path
    kind: str
    text: str


def should_skip(path: Path) -> bool:
    normalised = str(path)
    return any(part in path.parts or part in normalised for part in SKIP_PARTS)


def clean_line(line: str) -> str:
    return re.sub(r"\s+", " ", line).strip(" \t\r\n-")


def unique_lines(lines: list[str], limit: int = 80) -> list[str]:
    seen: set[str] = set()
    output: list[str] = []
    for line in lines:
        cleaned = clean_line(line)
        key = cleaned.lower()
        if len(cleaned) < 4 or key in seen or NOISE_LINE.match(cleaned):
            continue
        seen.add(key)
        output.append(cleaned)
        if len(output) >= limit:
            break
    return output


def public_label(text: str) -> str:
    return re.sub(r"Rob[ _-]+Johnston", "Learner", text, flags=re.I)


def is_safe_criteria_line(line: str) -> bool:
    if not line or len(line) > 260 or NOISE_LINE.match(line):
        return False
    if PRIVATE_OR_EVIDENCE_LINE.search(line):
        return False
    return True


def extract_criteria_lines(text: str, limit: int = 80) -> list[str]:
    raw_lines = [clean_line(line) for line in text.splitlines()]
    raw_lines = [line for line in raw_lines if line]
    output: list[str] = []

    for index, line in enumerate(raw_lines):
        if not is_safe_criteria_line(line):
            continue

        in_assessment_context = any(
            ASSESSMENT_CONTEXT.search(raw_lines[back_index])
            for back_index in range(max(0, index - 4), index)
        )
        starts_like_criterion = bool(CRITERION_START.match(line))
        has_criterion_verb = bool(CRITERION_VERB.search(line))

        if not starts_like_criterion or not (has_criterion_verb or in_assessment_context):
            continue

        parts = [line]
        lookahead = index + 1
        while lookahead < min(len(raw_lines), index + 5):
            next_line = raw_lines[lookahead]
            if CRITERION_START.match(next_line):
                break
            if not is_safe_criteria_line(next_line):
                lookahead += 1
                continue
            if next_line.startswith(("•", "-", "*")) or line.lower().endswith(("for", "of", "use", "using")):
                parts.append(next_line.lstrip("•-* "))
            lookahead += 1

        output.append(" ".join(parts))

    return unique_lines([public_label(line) for line in output], limit=limit)


def extract_text_file(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def extract_docx(path: Path) -> str:
    document = Document(path)
    parts = [paragraph.text for paragraph in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def extract_pptx(path: Path) -> str:
    deck = Presentation(path)
    parts: list[str] = []
    for slide_number, slide in enumerate(deck.slides, start=1):
        parts.append(f"Slide {slide_number}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return "\n".join(parts)


def extract_pdf(path: Path) -> str:
    if PdfReader is None:
        return ""
    reader = PdfReader(str(path))
    pages = []
    for page in reader.pages[:80]:
        pages.append(page.extract_text() or "")
    return "\n".join(pages)


EXTRACTORS: dict[str, Callable[[Path], str]] = {
    ".txt": extract_text_file,
    ".md": extract_text_file,
    ".ipynb": extract_text_file,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".pdf": extract_pdf,
}


def iter_source_paths(root: Path) -> list[Path]:
    paths: list[Path] = []
    for path in root.rglob("*"):
        if should_skip(path) or not path.is_file():
            continue
        if path.suffix.lower() not in SUPPORTED_EXTENSIONS:
            continue
        if not RELEVANT_PATH.search(str(path.relative_to(root))):
            continue
        paths.append(path)
    return sorted(paths, key=lambda item: str(item).lower())


def source_title(relative_path: str) -> str:
    name = Path(relative_path).stem
    name = re.sub(r"[_-]+", " ", name)
    return public_label(clean_line(name))


def source_ref(relative_path: str) -> str:
    return public_label(relative_path)


def source_unit(relative_path: str) -> str:
    match = re.search(r"(Unit\s*\d+[^\\/]*|unit-\d+[^\\/]*)", relative_path, re.I)
    if match:
        return clean_line(match.group(1)).title()
    return "AI Programming"


def source_id(relative_path: str, index: int) -> str:
    slug = re.sub(r"[^a-z0-9]+", "-", relative_path.lower()).strip("-")
    return f"src-{index + 1:03d}-{slug[:70]}"


def build_index(root: Path) -> dict[str, object]:
    source_paths = iter_source_paths(root)
    indexed_sources: list[dict[str, object]] = []
    skipped: list[dict[str, str]] = []

    for index, path in enumerate(source_paths):
        relative_path = path.relative_to(root).as_posix()
        extractor = EXTRACTORS[path.suffix.lower()]
        try:
            text = extractor(path)
        except Exception as exc:  # keep indexing resilient
            skipped.append({"path": relative_path, "reason": type(exc).__name__})
            continue

        criteria_lines = extract_criteria_lines(text, limit=80)
        if not criteria_lines:
            continue

        indexed_sources.append(
            {
                "id": source_id(source_ref(relative_path), index),
                "title": source_title(relative_path),
                "unit": source_unit(relative_path),
                "path": source_ref(relative_path),
                "type": path.suffix.lower().lstrip("."),
                "criteriaText": "\n".join(criteria_lines),
                "criteriaCount": len(criteria_lines),
                "keywords": sorted(set(re.findall(r"\b(?:AC\s*)?\d+\.\d+\b", "\n".join(criteria_lines), flags=re.I))),
            }
        )

    default_id = ""
    preferred_default = ("ac2.3 commit version control evidence", "ac2.3", "2.3 commit")
    for phrase in preferred_default:
        for source in indexed_sources:
            haystack = f"{source['title']} {source['path']} {source['criteriaText']}".lower()
            if phrase in haystack:
                default_id = str(source["id"])
                break
        if default_id:
            break

    for source in indexed_sources:
        if default_id:
            break
        haystack = f"{source['title']} {source['path']} {source['criteriaText']}".lower()
        if "2.3" in haystack or "ac2.3" in haystack:
            default_id = str(source["id"])
            break
    if not default_id and indexed_sources:
        default_id = str(indexed_sources[0]["id"])

    return {
        "generatedAt": datetime.now(timezone.utc).isoformat(),
        "rootLabel": "Local AI Programming source tree",
        "scannedSourceCount": len(source_paths),
        "indexedSourceCount": len(indexed_sources),
        "skippedSourceCount": len(skipped),
        "defaultSourceId": default_id,
        "sources": indexed_sources,
    }


def write_js_module(index: dict[str, object], out_path: Path) -> None:
    payload = json.dumps(index, indent=2)
    out_path.write_text(
        "// Generated by scripts/build_assessment_index.py. Do not edit by hand.\n"
        f"export const assessmentIndex = {payload};\n\n"
        "export const assessmentSources = assessmentIndex.sources;\n"
        "export const defaultAssessmentSourceId = assessmentIndex.defaultSourceId;\n",
        encoding="utf-8",
    )


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build the OG A.L.I. assessment criteria index")
    parser.add_argument(
        "--root",
        type=Path,
        default=Path(r"C:\OGBuild\College\Ai Programming"),
        help="Local AI Programming source folder",
    )
    parser.add_argument("--out", type=Path, default=Path("src/assessmentCriteria.js"))
    parser.add_argument("--summary", action="store_true", help="Print JSON summary after writing")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    root = args.root.resolve()
    if not root.exists():
        raise SystemExit(f"Source root not found: {root}")
    index = build_index(root)
    write_js_module(index, args.out)
    if args.summary:
        print(
            json.dumps(
                {
                    "scannedSourceCount": index["scannedSourceCount"],
                    "indexedSourceCount": index["indexedSourceCount"],
                    "skippedSourceCount": index["skippedSourceCount"],
                    "defaultSourceId": index["defaultSourceId"],
                },
                indent=2,
            )
        )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
